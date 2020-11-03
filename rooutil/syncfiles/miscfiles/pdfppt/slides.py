from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os, sys, commands


def addPic(slide, img, alignLeft=False, alignRight=False, withTextbox=False):
    global slideWidth, slideHeight

    im = Image.open(img)
    imgWidth, imgHeight = im.size
    imgRatio = 1.0*imgWidth/imgHeight


    # full
    top = 1.4 # under title
    height = slideHeight - top
    width = imgRatio*height
    left = (slideWidth - width)/2.0

    # left
    if(alignLeft):
        width = slideWidth / 2.0
        left = 0
        height = width / imgRatio
        top = 2.0

    # right
    if(alignRight):
        widthDivision = 2.0

        if(withTextbox): widthDivision = 1.5

        width = slideWidth / widthDivision
        left = slideWidth - width 
        height = width / imgRatio
        top = 2.0

        if(withTextbox): top = 1.0

    if(withTextbox):
        leftTB, topTB, widthTB, heightTB = map(Inches, [0.7,1.8,slideWidth-width-0.7,slideHeight-1.8])
        txBox = slide.shapes.add_textbox(leftTB, topTB, widthTB, heightTB)
        tf = txBox.textframe
        tf.word_wrap = True

        tf.text = "This is text inside a textbox"





    pic = slide.shapes.add_picture(img, int(Inches(left)), int(Inches(top)), int(Inches(width)), int(Inches(height)))



def picSlide(prs, img1):
    titleOnlyLayout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(titleOnlyLayout)

    addPic(slide, img1)

    slide.shapes.title.text = img1.split(".")[0]

def picSlideTwo(prs, img1, img2):
    titleOnlyLayout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(titleOnlyLayout)

    addPic(slide, img1, alignLeft=True)
    addPic(slide, img2, alignRight=True)

    slide.shapes.title.text = img1.split(".")[0] + "," + img2.split(".")[0]

def picSlideText(prs, img1):
    titleOnlyLayout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(titleOnlyLayout)

    addPic(slide, img1, alignRight=True, withTextbox=True)


    slide.shapes.title.text = img1.split(".")[0]

def pdfToPng(pdf):
    os.system("pdftopng %s" % pdf)
    return pdf.replace(".pdf",".png")



listFile = "list.txt"

if(len( sys.argv ) > 1):
    listFile = sys.argv[-1]


prs = Presentation(os.getenv('HOME')+'/syncfiles/miscfiles/pdfppt/template.pptx')

slideWidth, slideHeight = prs.slide_width / 914400.0, prs.slide_height / 914400.0


pdfSlides = []
if(not os.path.isfile(listFile)):
    pdfSlides = commands.getstatusoutput("ls -1 *.pdf")[1]
    pdfSlides = [pdf.strip() for pdf in pdfSlides.split("\n")]
    # out = open(listFile, 'w')
    # out.write(pdfFiles)
    # out.close()
    print pdfSlides
else:
    pdfSlides = open(listFile, 'r').read()
    pdfSlides = [pdf.strip() for pdf in pdfSlides.split("\n")]
    print pdfSlides

for s in pdfSlides:
    if(len(s) < 3): 
        continue

    elements = s.split()

    if(len(elements) < 2):
        print "full slide for", elements[0]
        picSlide(prs, pdfToPng(elements[0]))
    else:
        if("text" in elements[0]):
            print "textbox + half slide for", elements[1]
            picSlideText(prs, pdfToPng(elements[1]))
        else:
            print "split slide for", elements[0], "and", elements[1]
            picSlideTwo(prs, pdfToPng(elements[0]), pdfToPng(elements[1]))


prs.save('output.pptx')
# # print pdfFiles



